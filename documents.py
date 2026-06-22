"""Document text extraction for supported file types.

Extracts plain text suitable for injection into an LLM prompt context.
Imports are deferred so the app starts even if a library is missing for
a format the user never uses.
"""

from __future__ import annotations

import os
import platform
import sys

SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".csv")


def _tessdata_path() -> str | None:
    """Locate tessdata directory, preferring the bundled copy when frozen."""
    meipass = getattr(sys, "_MEIPASS", None)
    if meipass:
        p = os.path.join(meipass, "tessdata")
        return p if os.path.isdir(p) else None
    # Development / run-from-source: check common system install locations.
    plat = platform.system()
    if plat == "Linux":
        candidates = [
            "/usr/share/tesseract-ocr/5/tessdata",
            "/usr/share/tesseract-ocr/4.00/tessdata",
            "/usr/share/tessdata",
        ]
    elif plat == "Darwin":
        candidates = ["/opt/homebrew/share/tessdata", "/usr/local/share/tessdata"]
    elif plat == "Windows":
        candidates = [
            r"C:\Program Files\Tesseract-OCR\tessdata",
            r"C:\Program Files (x86)\Tesseract-OCR\tessdata",
        ]
    else:
        candidates = []
    return next((d for d in candidates if os.path.isdir(d)), None)


def _ensure_tesseract_on_path() -> None:
    """When frozen, prepend sys._MEIPASS to PATH so the bundled tesseract binary
    is found by PyMuPDF's OCR engine (which calls it as a subprocess)."""
    meipass = getattr(sys, "_MEIPASS", None)
    if meipass and meipass not in os.environ.get("PATH", ""):
        os.environ["PATH"] = meipass + os.pathsep + os.environ.get("PATH", "")

FILTER = (
    "Documents (*.pdf *.docx *.xlsx *.pptx *.txt *.md *.csv);;"
    "PDF (*.pdf);;"
    "Word (*.docx);;"
    "Excel (*.xlsx);;"
    "PowerPoint (*.pptx);;"
    "Text (*.txt *.md *.csv);;"
    "All files (*)"
)


def extract_text(path: str) -> str:
    """Return the textual content of *path* as a plain UTF-8 string."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _pdf(path)
    if ext == ".docx":
        return _docx(path)
    if ext == ".xlsx":
        return _xlsx(path)
    if ext == ".pptx":
        return _pptx(path)
    if ext in (".txt", ".md", ".csv"):
        return _plain(path)
    raise ValueError(f"Unsupported file type: {ext!r}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}")


def _pdf(path: str) -> str:
    import fitz  # PyMuPDF
    _ensure_tesseract_on_path()
    tessdata = _tessdata_path()
    doc = fitz.open(path)
    pages = []
    for page in doc:
        text = page.get_text().strip()
        # Fewer than 50 chars on a page almost certainly means it's a scanned image.
        if len(text) < 50 and tessdata:
            try:
                tp = page.get_textpage_ocr(tessdata=tessdata, language="eng", dpi=150)
                text = page.get_text(textpage=tp).strip()
            except Exception:
                pass  # OCR unavailable; keep whatever text was found
        pages.append(text)
    doc.close()
    return "\n\n".join(p for p in pages if p)


def _docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            lines.append(t)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"Sheet: {name}")
        for row in ws.iter_rows(values_only=True):
            line = "\t".join("" if v is None else str(v) for v in row)
            if line.strip():
                parts.append(line)
    wb.close()
    return "\n".join(parts)


def _pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            para.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            if para.text.strip()
        ]
        if texts:
            parts.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _plain(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()
